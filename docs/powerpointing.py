from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Define the content for each slide
slides_content = [
    ("The Detrimental Effects of Pornography", "Understanding the Impact on Individuals and Society"),
    ("Introduction", "Define pornography. Brief overview of the prevalence of pornography in today's society"),
    ("Scope of the Issue", "Statistics on the consumption of pornography globally. Age groups most affected"),
    ("Psychological Impact", "Addiction, Desensitization, Objectification"),
    ("Social Impact", "Relationship issues, Sexual attitudes and behaviors, Violence and aggression"),
    ("Physical Impact", "Erectile dysfunction, Brain changes, Health risks"),
    ("Ethical and Legal Considerations", "Exploitation, Legal issues, Child pornography"),
    ("Cultural and Societal Implications", "Shaping cultural norms, Gender dynamics, Impact on youth"),
    ("Addressing the Issue", "Education and awareness, Support and treatment, Advocacy and policy"),
    ("Conclusion", "Recap of key points, Call to action"),
    ("Additional Resources", "List of references, Hotlines and support organizations")
]

# Create slides with content
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use layout for title and content
    title_placeholder = slide.placeholders[0]
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Save the presentation
prs.save("Detrimental_Effects_of_Pornography.pptx")
